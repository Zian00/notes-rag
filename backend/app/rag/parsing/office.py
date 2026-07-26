import io
import re

from docx import Document as DocxDocument
from pptx import Presentation

from app.rag.types import ParsedDocument, Segment

# Matches Word paragraph style names like "Heading 1", "Heading 2" (case-insensitive,
# optional whitespace before the digit) to extract the heading level.
_HEADING_STYLE_RE = re.compile(r"heading\s*(\d+)", re.IGNORECASE)


class PptxParser:
    """One segment per slide; section = slide title, page_number = slide number."""

    def parse(self, data: bytes, content_type: str) -> ParsedDocument:
        prs = Presentation(io.BytesIO(data))
        segments: list[Segment] = []
        # One slide → one segment. enumerate(..., start=1) gives a 1-based slide number we
        # store as page_number (so retrieval can cite "slide 3").
        for index, slide in enumerate(prs.slides, start=1):
            # Read the title via slide.shapes.title.text directly. (Comparing shapes with
            # `shape == slide.shapes.title` does NOT work — python-pptx hands back fresh
            # proxy objects each access, so the identity check is always False.)
            title_shape = slide.shapes.title
            title = (
                title_shape.text
                if title_shape is not None and title_shape.text
                else None
            )
            # Collect text from every shape that has a text frame (title + body + textboxes).
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text
            ]
            segments.append(
                Segment(text="\n".join(texts), page_number=index, section=title)
            )
        return ParsedDocument(segments=segments, page_count=len(segments))


class DocxParser:
    """Segments split on heading-styled paragraphs, with a '>'-joined breadcrumb
    `section` for nested heading levels (mirrors heading_split.py's approach, but
    the level signal here is the paragraph's Word STYLE, e.g. 'Heading 2', not a
    '#' count)."""

    def parse(self, data: bytes, content_type: str) -> ParsedDocument:
        doc = DocxDocument(io.BytesIO(data))
        segments: list[Segment] = []
        # Stack of (level, heading_text). A new heading pops any stack entries at
        # the same or deeper level before being pushed, so the stack always reflects
        # the current breadcrumb path (e.g. [(1, "Lecture 4"), (2, "Neural Networks")]).
        heading_stack: list[tuple[int, str]] = []
        buffer: list[str] = []

        def breadcrumb() -> str | None:
            return " > ".join(h[1] for h in heading_stack) or None

        def flush() -> None:
            body = "\n".join(buffer).strip()
            if body or heading_stack:
                segments.append(Segment(text=body, section=breadcrumb()))

        # Same flush-on-heading pattern as Markdown, but here "heading" is a Word paragraph
        # STYLE (e.g. "Heading 1"), not a '#'. Word has no fixed page count, so page_count=None.
        for para in doc.paragraphs:
            style = (para.style.name or "") if para.style else ""
            match = _HEADING_STYLE_RE.match(style.strip())
            if match and para.text.strip():
                flush()
                buffer = []
                level = int(match.group(1))
                while heading_stack and heading_stack[-1][0] >= level:
                    heading_stack.pop()
                heading_stack.append((level, para.text.strip()))
            elif para.text.strip():
                buffer.append(para.text)
        flush()
        if not segments:
            segments = [Segment(text="")]  # guarantee at least one segment for an empty doc
        return ParsedDocument(segments=segments, page_count=None)
