import io
import shutil

import pytest
from app.rag.parsing import ParserDispatcher
from app.rag.parsing.image import ImageParser
from app.rag.parsing.office import DocxParser, PptxParser
from app.rag.parsing.pdf import PdfParser
from app.rag.parsing.text import TextParser
from docx import Document as DocxDocument
from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from tests.fakes import FakeOcrProvider


def test_text_parser_plain_is_single_segment():
    doc = TextParser().parse(b"line one\nline two", "text/plain")
    assert len(doc.segments) == 1
    assert "line one" in doc.segments[0].text


def test_text_parser_markdown_splits_on_headings():
    md = b"# Intro\nhello\n## Details\nworld"
    doc = TextParser().parse(md, "text/markdown")
    sections = [s.section for s in doc.segments]
    assert "Intro" in sections
    assert any(s and "Details" in s for s in sections)


def test_markdown_nested_headings_produce_breadcrumb_sections():
    md = "# Lecture 4\nIntro text.\n## Neural Networks\nBody about neurons.\n"
    parsed = TextParser().parse(md.encode(), "text/markdown")
    sections = [s.section for s in parsed.segments]
    assert "Lecture 4" in sections
    assert "Lecture 4 > Neural Networks" in sections


def test_pptx_parser_one_segment_per_slide():
    prs = Presentation()
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {i}"
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = f"body {i}"
    buf = io.BytesIO()
    prs.save(buf)
    doc = PptxParser().parse(buf.getvalue(), "application/...")
    assert len(doc.segments) == 2
    assert doc.segments[0].page_number == 1
    assert doc.segments[0].section == "Slide 0"
    assert "body 0" in doc.segments[0].text


def test_docx_parser_extracts_text():
    d = DocxDocument()
    d.add_paragraph("first para")
    d.add_paragraph("second para")
    buf = io.BytesIO()
    d.save(buf)
    doc = DocxParser().parse(buf.getvalue(), "application/...")
    joined = " ".join(s.text for s in doc.segments)
    assert "first para" in joined and "second para" in joined


def test_image_parser_uses_ocr():
    img = Image.new("RGB", (40, 20), "white")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    doc = ImageParser(FakeOcrProvider("scanned words")).parse(buf.getvalue(), "image/png")
    assert len(doc.segments) == 1
    assert doc.segments[0].text == "scanned words"
    assert doc.segments[0].page_number == 1


def test_pdf_parser_extracts_text_layer():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=14)
    pdf.cell(text="hello pdf world")
    data = bytes(pdf.output())
    doc = PdfParser(ocr=FakeOcrProvider("FALLBACK"), ocr_enabled=True, min_chars=5).parse(
        data, "application/pdf"
    )
    text = " ".join(s.text for s in doc.segments)
    assert "hello pdf world" in text
    assert "FALLBACK" not in text  # text layer present -> no OCR
    assert doc.page_count == 1


def test_pdf_parser_falls_back_to_ocr_for_image_only_page():
    if shutil.which("pdftoppm") is None:
        pytest.skip("poppler (pdftoppm) not installed")
    # An image saved as a 1-page PDF has no text layer -> OCR fallback fires.
    img = Image.new("RGB", (200, 80), "white")
    buf = io.BytesIO()
    img.save(buf, format="PDF")
    doc = PdfParser(ocr=FakeOcrProvider("FALLBACK TEXT"), ocr_enabled=True, min_chars=5).parse(
        buf.getvalue(), "application/pdf"
    )
    text = " ".join(s.text for s in doc.segments)
    assert "FALLBACK TEXT" in text


def test_dispatcher_routes_by_content_type():
    dispatcher = ParserDispatcher(ocr=FakeOcrProvider(), ocr_enabled=True, min_chars=5)
    doc = dispatcher.parse(b"plain words", "text/plain")
    assert "plain words" in doc.segments[0].text
